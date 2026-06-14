"""
Build TCL 2026 Training Presentation — 72 slides
Apple Keynote style: black background, TCL red, giant white typography
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xE8, 0x00, 0x2D)
GRAY   = RGBColor(0x88, 0x88, 0x88)
LGRAY  = RGBColor(0xBB, 0xBB, 0xBB)
DGRAY  = RGBColor(0x22, 0x22, 0x22)
GOLD   = RGBColor(0xD4, 0xAF, 0x37)

W  = Inches(13.33)   # 16:9 widescreen width
H  = Inches(7.5)     # 16:9 widescreen height

OUTPUT = "/root/.claude/TCL_2026_Training.pptx"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def add_text(slide, text, x, y, w, h,
             size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Arial"
    return txBox


def add_line(slide, x, y, w, color=RED, thickness=Pt(3)):
    """Horizontal divider line."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_feat_row(slide, label, value, y):
    """One row in a feature table."""
    cx = Inches(2.5)
    cw = Inches(8.33)
    add_text(slide, label, cx, y, Inches(4), Inches(0.4),
             size=20, bold=False, color=LGRAY, align=PP_ALIGN.LEFT)
    add_text(slide, value, cx + Inches(4), y, Inches(4.33), Inches(0.4),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    # thin separator
    sep = slide.shapes.add_shape(1, cx, y + Inches(0.38), cw, Inches(0.01))
    sep.fill.solid(); sep.fill.fore_color.rgb = DGRAY
    sep.line.fill.background()


def label_slide(slide, text, color=RED):
    """Small uppercase label at top center."""
    add_text(slide, text.upper(), Inches(0), Inches(0.5), W, Inches(0.4),
             size=14, bold=True, color=color, align=PP_ALIGN.CENTER)


def hero_number(slide, number_str, spec_str="", model_str="", num_size=160):
    """Giant centered number with optional spec and model below."""
    add_text(slide, number_str,
             Inches(0), Inches(1.2), W, Inches(4.5),
             size=num_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if spec_str:
        add_text(slide, spec_str,
                 Inches(0), Inches(5.2), W, Inches(0.9),
                 size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    if model_str:
        add_text(slide, model_str,
                 Inches(0), Inches(6.0), W, Inches(0.7),
                 size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)


def divider(slide, y=Inches(3.8)):
    add_line(slide, Inches(6.17), y, Inches(1), RED)


# ══════════════════════════════════════════════════════════════════════════════
# PART 1 — OPENING
# ══════════════════════════════════════════════════════════════════════════════

# Slide 1 — TCL logo
s = new_slide()
add_text(s, "TCL", Inches(0), Inches(2.8), W, Inches(1.5),
         size=72, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Slide 2 — TCL + 2026
s = new_slide()
add_text(s, "TCL", Inches(0), Inches(1.2), W, Inches(1.2),
         size=48, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, "2026", Inches(0), Inches(2.4), W, Inches(4),
         size=160, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Slide 3 — Opening statement
s = new_slide()
label_slide(s, "Training Material")
add_text(s, "The Complete\nTCL TV Lineup.", Inches(0), Inches(1.6), W, Inches(3.2),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(4.7))
add_text(s, "10 Models. One Vision.", Inches(0), Inches(4.9), W, Inches(0.8),
         size=26, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 4 — Not a catalogue
s = new_slide()
add_text(s, "This is not a product catalogue.", Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.2),
         size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "This is your competitive edge.", Inches(0.5), Inches(3.8), Inches(12.33), Inches(1),
         size=36, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 5 — Let's begin
s = new_slide()
add_text(s, "Let's begin.", Inches(0), Inches(2.5), W, Inches(2),
         size=90, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(4.4))

# ══════════════════════════════════════════════════════════════════════════════
# PART 2 — OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════

# Slide 6 — What changed
s = new_slide()
label_slide(s, "2026")
add_text(s, "What changed\nthis year?", Inches(0), Inches(1.8), W, Inches(3.5),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Slide 7 — Everything.
s = new_slide()
add_text(s, "Everything.", Inches(0), Inches(1.5), W, Inches(4.5),
         size=130, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Slide 8 — Everything expanded
s = new_slide()
add_text(s, "Everything.", Inches(0), Inches(0.8), W, Inches(1.6),
         size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(2.3))
add_text(s, "SQD Technology — across the premium range.",
         Inches(1), Inches(2.6), Inches(11.33), Inches(0.7),
         size=24, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, "RGB MiniLED — for the ultra-large category.",
         Inches(1), Inches(3.3), Inches(11.33), Inches(0.7),
         size=24, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, "Dolby Vision 2 Max — exclusive to TCL.",
         Inches(1), Inches(4.0), Inches(11.33), Inches(0.7),
         size=24, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 9 — Full lineup table
s = new_slide()
label_slide(s, "2026 Full Lineup")
rows = [
    ("X11L",          "Flagship · SQD MiniLED · 10,000 nits"),
    ("RM9L",          "Ultra-Large · RGB MiniLED · 9,000 nits"),
    ("C8L",           "Premium · SQD MiniLED · 6,000 nits"),
    ("C7L",           "Best Seller · SQD MiniLED · 3,000 nits"),
    ("RM7L",          "Large · RGB MiniLED · 2,000 nits"),
    ("C6L",           "Mid Premium · QLED"),
    ("P8L / P7L / P6L", "Value · 4K Smart TV"),
    ("V6D / T6D",     "Entry · Smart TV"),
]
for i, (lbl, val) in enumerate(rows):
    add_feat_row(s, lbl, val, Inches(1.1) + Inches(0.72) * i)

# Slide 10 — Know the lineup
s = new_slide()
add_text(s, "Know the lineup.", Inches(0), Inches(2.0), W, Inches(1.5),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Own the floor.", Inches(0), Inches(3.5), W, Inches(1.5),
         size=72, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# PART 3 — X11L
# ══════════════════════════════════════════════════════════════════════════════

# Slide 11
s = new_slide()
label_slide(s, "Flagship")
add_text(s, "X11L", Inches(0), Inches(1.0), W, Inches(4.5),
         size=180, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "The pinnacle of display technology.",
         Inches(0), Inches(5.5), W, Inches(0.8),
         size=22, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

# Slide 12
s = new_slide()
add_text(s, "There is no TV\nabove the X11L.", Inches(0), Inches(1.8), W, Inches(4),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Slides 13–14: 10,000 nits
s = new_slide()
hero_number(s, "10,000", num_size=170)

s = new_slide()
hero_number(s, "10,000", "Peak Nits of Brightness", "TCL X11L", num_size=150)

# Slides 15–16: 20,000+ zones
s = new_slide()
hero_number(s, "20,000+", num_size=140)

s = new_slide()
hero_number(s, "20,000+", "Local Dimming Zones", "The most precise backlight ever built.", num_size=120)

# Slide 17: WHVA
s = new_slide()
label_slide(s, "Panel Technology")
add_text(s, "WHVA 2.0 Ultra", Inches(0), Inches(1.5), W, Inches(2),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(3.4))
add_text(s, "Wide High Viewing Angle — 2nd Generation.",
         Inches(0), Inches(3.6), W, Inches(0.8),
         size=26, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, "Perfect color from every seat in the room.",
         Inches(0), Inches(4.4), W, Inches(0.8),
         size=26, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 18: X11L specs table
s = new_slide()
label_slide(s, "X11L Full Specs")
specs = [
    ("Panel",           "WHVA 2.0 Ultra"),
    ("Backlight",       "SQD MiniLED"),
    ("Peak Brightness", "10,000 nits"),
    ("Dimming Zones",   "20,000+"),
    ("HDR Support",     "Dolby Vision · HDR10+"),
    ("Screen Sizes",    "75\"  ·  85\"  ·  98\""),
]
for i, (l, v) in enumerate(specs):
    add_feat_row(s, l, v, Inches(1.1) + Inches(0.8) * i)

# Slide 19: Sizes
s = new_slide()
label_slide(s, "Available In")
add_text(s, "75\"     85\"     98\"", Inches(0), Inches(2.0), W, Inches(2.5),
         size=90, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(4.4))
add_text(s, "From premium to palace.", Inches(0), Inches(4.7), W, Inches(0.8),
         size=28, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 20: X11L close
s = new_slide()
add_text(s, "There's nothing\nabove it.", Inches(0), Inches(1.8), W, Inches(3.5),
         size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "TCL X11L", Inches(0), Inches(5.4), W, Inches(0.7),
         size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# PART 4 — RM9L
# ══════════════════════════════════════════════════════════════════════════════

# Slide 21
s = new_slide()
label_slide(s, "Ultra-Large Premium")
add_text(s, "RM9L", Inches(0), Inches(1.0), W, Inches(4.5),
         size=180, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "RGB MiniLED. Reimagined.",
         Inches(0), Inches(5.5), W, Inches(0.8),
         size=22, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

# Slide 22: fragments feel
s = new_slide()
add_text(s, "Massive screen.", Inches(0), Inches(1.4), W, Inches(1.4),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Massive brightness.", Inches(0), Inches(2.8), W, Inches(1.4),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Massive impact.", Inches(0), Inches(4.2), W, Inches(1.4),
         size=72, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Slides 23–24: 9,000 nits
s = new_slide(); hero_number(s, "9,000", num_size=170)
s = new_slide(); hero_number(s, "9,000", "Peak Nits · RGB MiniLED", "TCL RM9L", num_size=150)

# Slide 25: Sizes
s = new_slide()
label_slide(s, "Screen Sizes")
add_text(s, "85\"     98\"     115\"", Inches(0), Inches(2.0), W, Inches(2.5),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(4.4))
add_text(s, "When bigger is the point.", Inches(0), Inches(4.7), W, Inches(0.8),
         size=28, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 26: RM9L specs
s = new_slide()
label_slide(s, "RM9L Full Specs")
specs = [
    ("Backlight",       "RGB MiniLED"),
    ("Peak Brightness", "9,000 nits"),
    ("Color Volume",    "Ultra-wide color gamut"),
    ("HDR Support",     "Dolby Vision · HDR10+"),
    ("Screen Sizes",    "85\"  ·  98\"  ·  115\""),
]
for i, (l, v) in enumerate(specs):
    add_feat_row(s, l, v, Inches(1.4) + Inches(0.85) * i)

# Slide 27: RM9L close
s = new_slide()
add_text(s, "For the customer who wants it all\n— and larger.",
         Inches(0.5), Inches(1.8), Inches(12.33), Inches(4),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# PART 5 — C8L
# ══════════════════════════════════════════════════════════════════════════════

# Slide 28
s = new_slide()
label_slide(s, "Premium")
add_text(s, "C8L", Inches(0), Inches(1.0), W, Inches(4.5),
         size=200, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "SQD MiniLED. Perfected.",
         Inches(0), Inches(5.5), W, Inches(0.8),
         size=22, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

# Slides 29–30: 6,000 nits
s = new_slide(); hero_number(s, "6,000", num_size=170)
s = new_slide(); hero_number(s, "6,000", "Peak Nits · SQD MiniLED", "TCL C8L", num_size=150)

# Slides 31–32: 4,032 zones
s = new_slide(); hero_number(s, "4,032", num_size=170)
s = new_slide(); hero_number(s, "4,032", "Local Dimming Zones", "Precision you can see.", num_size=150)

# Slide 33: Dolby Vision 2 Max EXCLUSIVE
s = new_slide()
label_slide(s, "World Exclusive")
add_text(s, "Dolby Vision 2 Max", Inches(0), Inches(1.4), W, Inches(2),
         size=76, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# Red badge
badge = s.shapes.add_shape(1, Inches(9.5), Inches(1.6), Inches(2.8), Inches(0.55))
badge.fill.solid(); badge.fill.fore_color.rgb = RED
badge.line.fill.background()
add_text(s, "EXCLUSIVE", Inches(9.5), Inches(1.6), Inches(2.8), Inches(0.55),
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(3.3))
add_text(s, "Only available on the TCL C8L.",
         Inches(0), Inches(3.6), W, Inches(0.8),
         size=28, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, "No other brand has this. Use it.",
         Inches(0), Inches(4.4), W, Inches(0.8),
         size=28, bold=False, color=RED, align=PP_ALIGN.CENTER)

# Slide 34: C8L specs
s = new_slide()
label_slide(s, "C8L Full Specs")
specs = [
    ("Technology",      "SQD MiniLED"),
    ("Peak Brightness", "6,000 nits"),
    ("Dimming Zones",   "4,032"),
    ("Exclusive HDR",   "★  Dolby Vision 2 Max"),
    ("HDR also",        "HDR10+"),
    ("Screen Sizes",    "55\"  to  98\""),
]
for i, (l, v) in enumerate(specs):
    add_feat_row(s, l, v, Inches(1.1) + Inches(0.8) * i)

# Slide 35: C8L close
s = new_slide()
add_text(s, "Premium picture.", Inches(0), Inches(1.8), W, Inches(1.6),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Exclusive technology.", Inches(0), Inches(3.4), W, Inches(1.6),
         size=80, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, "TCL C8L", Inches(0), Inches(5.4), W, Inches(0.7),
         size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# PART 6 — C7L BEST SELLER
# ══════════════════════════════════════════════════════════════════════════════

# Slide 36
s = new_slide()
label_slide(s, "Best Seller")
add_text(s, "C7L", Inches(0), Inches(1.0), W, Inches(4.5),
         size=200, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "SQD MiniLED. For everyone.",
         Inches(0), Inches(5.5), W, Inches(0.8),
         size=22, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

# Slide 37
s = new_slide()
add_text(s, "The one they'll\nask for by name.", Inches(0), Inches(1.8), W, Inches(3.8),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Slides 38–39: 3,000 nits
s = new_slide(); hero_number(s, "3,000", num_size=170)
s = new_slide(); hero_number(s, "3,000", "Peak Nits · SQD MiniLED", "TCL C7L", num_size=150)

# Slides 40–41: 2,176 zones
s = new_slide(); hero_number(s, "2,176", num_size=170)
s = new_slide(); hero_number(s, "2,176", "Local Dimming Zones",
                              "Blacks this deep have never been this affordable.", num_size=150)

# Slides 42–43: 144Hz
s = new_slide(); hero_number(s, "144Hz", num_size=150)
s = new_slide(); hero_number(s, "144Hz", "Refresh Rate", "Zero blur. Zero ghosting.", num_size=130)

# Slide 44: Use cases
s = new_slide()
add_text(s, "Gaming.", Inches(0), Inches(0.9), W, Inches(1.4),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Movies.", Inches(0), Inches(2.3), W, Inches(1.4),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Sports.", Inches(0), Inches(3.7), W, Inches(1.4),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Everything.", Inches(0), Inches(5.1), W, Inches(1.4),
         size=80, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Slide 45: C7L specs
s = new_slide()
label_slide(s, "C7L Full Specs")
specs = [
    ("Technology",      "SQD MiniLED"),
    ("Peak Brightness", "3,000 nits"),
    ("Dimming Zones",   "Up to 2,176"),
    ("Refresh Rate",    "144Hz"),
    ("HDR Support",     "Dolby Vision · HDR10+"),
    ("Screen Sizes",    "55\"  to  98\""),
    ("Color Combos",    "1.07 Billion"),
]
for i, (l, v) in enumerate(specs):
    add_feat_row(s, l, v, Inches(0.9) + Inches(0.72) * i)

# Slide 46: C7L close
s = new_slide()
add_text(s, "The best seller\nfor a reason.", Inches(0), Inches(1.8), W, Inches(3.5),
         size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "TCL C7L", Inches(0), Inches(5.4), W, Inches(0.7),
         size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# PART 7 — RM7L
# ══════════════════════════════════════════════════════════════════════════════

# Slide 47
s = new_slide()
label_slide(s, "Large Screen Value")
add_text(s, "RM7L", Inches(0), Inches(1.0), W, Inches(4.5),
         size=180, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "RGB MiniLED. Bold colors at scale.",
         Inches(0), Inches(5.5), W, Inches(0.8),
         size=22, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

# Slides 48–49: 2,000 nits
s = new_slide(); hero_number(s, "2,000", num_size=170)
s = new_slide(); hero_number(s, "2,000", "Peak Nits · RGB MiniLED", "TCL RM7L", num_size=150)

# Slide 50: RM7L close
s = new_slide()
add_text(s, "Big screen.", Inches(0), Inches(1.2), W, Inches(1.5),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Bold colors.", Inches(0), Inches(2.7), W, Inches(1.5),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Better value.", Inches(0), Inches(4.2), W, Inches(1.5),
         size=80, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# PART 8 — MID RANGE
# ══════════════════════════════════════════════════════════════════════════════

# Slide 51
s = new_slide()
label_slide(s, "Mid Range")
add_text(s, "C6L · P8L · P7L · P6L", Inches(0), Inches(2.0), W, Inches(2),
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(4.0))
add_text(s, "Smart. Sharp. Affordable.",
         Inches(0), Inches(4.3), W, Inches(0.8),
         size=28, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 52: C6L
s = new_slide()
label_slide(s, "QLED · Mid Premium")
add_text(s, "C6L", Inches(0), Inches(1.0), W, Inches(3.5),
         size=160, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(4.4))
specs = [("Technology", "QLED"), ("Resolution", "4K Ultra HD"), ("HDR", "Dolby Vision · HDR10+")]
for i, (l, v) in enumerate(specs):
    add_feat_row(s, l, v, Inches(4.8) + Inches(0.7) * i)

# Slide 53: P8L
s = new_slide()
label_slide(s, "MiniLED Entry · Value Premium")
add_text(s, "P8L", Inches(0), Inches(1.0), W, Inches(3.5),
         size=160, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(4.4))
specs = [("Technology", "MiniLED"), ("Resolution", "4K Ultra HD"), ("HDR", "HDR10+")]
for i, (l, v) in enumerate(specs):
    add_feat_row(s, l, v, Inches(4.8) + Inches(0.7) * i)

# Slide 54: P7L + P6L
s = new_slide()
label_slide(s, "4K Smart TV · Value")
# Two columns
add_text(s, "P7L", Inches(1.0), Inches(1.5), Inches(5.17), Inches(2),
         size=110, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "4K · HDR10 · Smart TV\nGreat everyday viewing\nfor any family.",
         Inches(1.0), Inches(3.8), Inches(5.17), Inches(2),
         size=22, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)
# Vertical line
vl = s.shapes.add_shape(1, Inches(6.62), Inches(1.5), Inches(0.02), Inches(4.5))
vl.fill.solid(); vl.fill.fore_color.rgb = DGRAY; vl.line.fill.background()
add_text(s, "P6L", Inches(7.17), Inches(1.5), Inches(5.17), Inches(2),
         size=110, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "4K · HDR · Smart TV\nThe entry point\nto the TCL family.",
         Inches(7.17), Inches(3.8), Inches(5.17), Inches(2),
         size=22, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 55: Mid-range close
s = new_slide()
add_text(s, "A TCL for\nevery budget.", Inches(0), Inches(1.5), W, Inches(3.5),
         size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(5.0))
add_text(s, "No customer walks away empty-handed.",
         Inches(0), Inches(5.2), W, Inches(0.8),
         size=26, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# PART 9 — ENTRY LEVEL
# ══════════════════════════════════════════════════════════════════════════════

# Slide 56
s = new_slide()
label_slide(s, "Entry Level")
add_text(s, "V6D & T6D", Inches(0), Inches(1.8), W, Inches(2.2),
         size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(4.0))
add_text(s, "Smart TV. Real quality. TCL price.",
         Inches(0), Inches(4.3), W, Inches(0.8),
         size=28, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 57: Entry specs
s = new_slide()
label_slide(s, "V6D & T6D")
specs = [
    ("Resolution",   "Full HD / 4K"),
    ("Smart OS",     "Google TV"),
    ("HDR",          "HDR10"),
    ("Connectivity", "WiFi · Bluetooth · HDMI"),
    ("Positioning",  "Best first TV"),
]
for i, (l, v) in enumerate(specs):
    add_feat_row(s, l, v, Inches(1.5) + Inches(0.85) * i)

# Slide 58
s = new_slide()
add_text(s, "Everyone deserves", Inches(0), Inches(2.0), W, Inches(1.6),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "a great TV.", Inches(0), Inches(3.6), W, Inches(1.6),
         size=72, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# PART 10 — SQD TECHNOLOGY DEEP DIVE
# ══════════════════════════════════════════════════════════════════════════════

# Slide 59
s = new_slide()
label_slide(s, "Technology Deep Dive")
add_text(s, "SQD", Inches(0), Inches(0.8), W, Inches(5.5),
         size=200, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Slide 60
s = new_slide()
add_text(s, "Sub-pixel\nQuantum Dot.", Inches(0), Inches(1.2), W, Inches(3.5),
         size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(4.8))
add_text(s, "The next evolution beyond QLED.",
         Inches(0), Inches(5.0), W, Inches(0.7),
         size=24, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, "TCL-exclusive. Patent-protected.",
         Inches(0), Inches(5.7), W, Inches(0.7),
         size=24, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 61: Regular QLED vs SQD
s = new_slide()
add_text(s, "Regular QLED:", Inches(0), Inches(1.8), W, Inches(1.2),
         size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "One filter. Whole screen.", Inches(0), Inches(3.0), W, Inches(1.2),
         size=52, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

# Slide 62: SQD
s = new_slide()
add_text(s, "SQD:", Inches(0), Inches(1.4), W, Inches(1.2),
         size=72, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, "Every pixel.", Inches(0), Inches(2.6), W, Inches(1.4),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Its own filter.", Inches(0), Inches(4.0), W, Inches(1.4),
         size=72, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slides 63–64: 1.07 Billion
s = new_slide(); hero_number(s, "1.07B", num_size=160)
s = new_slide(); hero_number(s, "1.07B", "Color Combinations",
                              "The human eye can see about 10 million.", num_size=130)

# Slide 65: SQD vs competition
s = new_slide()
label_slide(s, "SQD vs The Competition")
# Three columns
cols = [
    ("QLED",  GRAY,  "Panel-level quantum dot filter.\nGood color. Limited precision.\nOne filter for millions of pixels."),
    ("OLED",  GRAY,  "Self-emissive pixels.\nGreat black levels.\nLower brightness. Burn-in risk."),
    ("SQD ★", RED,   "Sub-pixel level filter.\nPeak brightness of MiniLED.\nColor precision per pixel."),
]
for i, (title, col, desc) in enumerate(cols):
    x = Inches(0.5) + Inches(4.28) * i
    add_text(s, title, x, Inches(1.4), Inches(3.8), Inches(1),
             size=44, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, desc, x, Inches(2.6), Inches(3.8), Inches(2.5),
             size=20, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)
    if i < 2:
        vl = s.shapes.add_shape(1, x + Inches(3.9), Inches(1.4), Inches(0.02), Inches(4))
        vl.fill.solid(); vl.fill.fore_color.rgb = DGRAY; vl.line.fill.background()

# ══════════════════════════════════════════════════════════════════════════════
# PART 11 — SALES TRAINING
# ══════════════════════════════════════════════════════════════════════════════

# Slide 66
s = new_slide()
label_slide(s, "Sales Training")
add_text(s, "How to sell TCL\nin 2026.", Inches(0), Inches(1.8), W, Inches(3.5),
         size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Slide 67: 3-step close
s = new_slide()
label_slide(s, "The 3-Step Close")
steps = [
    "01   Ask: \"What do you watch most?\"",
    "02   Match: Show the right model.",
    "03   Demo: Let the screen close the sale.",
]
for i, st in enumerate(steps):
    add_text(s, st, Inches(1.5), Inches(1.8) + Inches(1.5) * i, Inches(10.33), Inches(1.2),
             size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Slide 68: Show don't tell
s = new_slide()
add_text(s, "Show.", Inches(0), Inches(1.8), W, Inches(1.6),
         size=100, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Don't tell.", Inches(0), Inches(3.4), W, Inches(1.6),
         size=100, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, "The SQD difference is visible in 5 seconds.",
         Inches(0), Inches(5.4), W, Inches(0.8),
         size=24, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)

# Slide 69: In-store demo script
s = new_slide()
label_slide(s, "The SQD Demo — In-Store Script")
script = [
    "1.  Point to any QLED TV next to C7L / C8L.",
    "2.  Say: \"See the green on that screen?\"",
    "3.  Point to TCL: \"Now look at ours.\"",
    "4.  Say: \"That's SQD — sub-pixel filtering.\"",
    "5.  Ask: \"Can you see the difference?\"",
]
for i, line in enumerate(script):
    add_text(s, line, Inches(1.5), Inches(1.5) + Inches(1.0) * i, Inches(10.33), Inches(0.9),
             size=26, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

# Slide 70: Upsell guide
s = new_slide()
label_slide(s, "Upsell Path")
upsell = [
    ("Budget-conscious",       "C7L — Best value SQD"),
    ("Wants the best picture", "C8L — Dolby Vision 2 Max"),
    ("Wants extreme size",     "RM9L — 85\" to 115\""),
    ("Money is no object",     "X11L — The pinnacle"),
    ("Gaming focus",           "C7L — 144Hz SQD"),
    ("First-time buyer",       "P6L or P7L — Smart 4K"),
]
for i, (l, v) in enumerate(upsell):
    add_feat_row(s, l, v, Inches(1.1) + Inches(0.8) * i)

# ══════════════════════════════════════════════════════════════════════════════
# PART 12 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════

# Slide 71
s = new_slide()
add_text(s, "You are now\nthe expert.", Inches(0), Inches(1.8), W, Inches(3.5),
         size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Slide 72: TCL CERTIFIED
s = new_slide()
add_text(s, "ACHIEVEMENT", Inches(0), Inches(0.8), W, Inches(0.6),
         size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "TCL\nCERTIFIED.", Inches(0), Inches(1.2), W, Inches(4),
         size=130, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_line(s, Inches(6.17), Inches(5.2), Inches(1), GOLD)
add_text(s, "You now sell the best TV technology in the world.",
         Inches(0), Inches(5.5), W, Inches(0.8),
         size=24, bold=False, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s, "TCL  ·  2026", Inches(0), Inches(6.5), W, Inches(0.7),
         size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"✅  Saved: {OUTPUT}")
print(f"   Slides: {len(prs.slides)}")

import os
print(f"   Size:   {os.path.getsize(OUTPUT):,} bytes")
