"""
TCL 2026 Training PPTX v2
- Pillow-generated cinematic dark backgrounds (5 types)
- PowerPoint Morph transitions (number morphs as spec appears)
- Click-triggered Fade-in animations on supporting text
"""

import numpy as np
from PIL import Image
import io, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
BLACK = RGBColor(0,    0,    0   )
WHITE = RGBColor(255,  255,  255 )
RED   = RGBColor(0xE8, 0x00, 0x2D)
GRAY  = RGBColor(0x88, 0x88, 0x88)
LGRAY = RGBColor(0xBB, 0xBB, 0xBB)
DGRAY = RGBColor(0x22, 0x22, 0x22)
GOLD  = RGBColor(0xD4, 0xAF, 0x37)

SLD_W  = Inches(13.33)
SLD_H  = Inches(7.5)
OUTPUT = "/root/.claude/TCL_2026_Training_v2.pptx"

prs = Presentation()
prs.slide_width  = SLD_W
prs.slide_height = SLD_H
blank = prs.slide_layouts[6]

# ── Background generators ─────────────────────────────────────────────────────
IW, IH = 1920, 1080

def _g():
    return np.mgrid[0:IH, 0:IW].astype(np.float32)

def bg_red_glow():
    y, x = _g()
    d = np.sqrt((x - IW*0.78)**2 + (y - IH*0.82)**2)
    g = np.clip(1 - d/(IW*0.42), 0, 1)**2 * 0.22
    r = (g*232).astype(np.uint8)
    b = (g*30 ).astype(np.uint8)
    return Image.fromarray(np.stack([r, np.zeros_like(r), b], 2))

def bg_dark_vignette(tint=(8, 8, 12)):
    y, x = _g()
    d = np.sqrt((x - IW/2)**2 + (y - IH/2)**2)
    v = np.clip(1 - d/(IW*0.72), 0, 1)
    arr = np.zeros((IH, IW, 3), dtype=np.uint8)
    for c, (base, add) in enumerate(zip(tint, (5, 5, 7))):
        arr[:,:,c] = np.clip(base + (v*add).astype(np.uint8), 0, 255)
    return Image.fromarray(arr)

def bg_number_glow(color=(30, 20, 20)):
    y, x = _g()
    d = np.sqrt((x - IW*0.5)**2 + (y - IH*0.5)**2)
    g = np.clip(1 - d/(IW*0.55), 0, 1)**3 * 0.18
    arr = np.zeros((IH, IW, 3), dtype=np.uint8)
    for c, cv in enumerate(color):
        arr[:,:,c] = (g*cv).astype(np.uint8)
    return Image.fromarray(arr)

def bg_tech_grid():
    arr = np.zeros((IH, IW, 3), dtype=np.uint8)
    sp = 56
    for gy in range(0, IH+sp, sp):
        for gx in range(0, IW+sp, sp):
            off = (sp//2) if ((gy//sp) % 2) else 0
            cx_ = gx + off
            for dy in range(-1, 2):
                for dx in range(-1, 2):
                    ry, rx = gy+dy, cx_+dx
                    if 0 <= ry < IH and 0 <= rx < IW:
                        arr[ry, rx] = [22, 4, 7]
    y, x = _g()
    d = np.sqrt((x - IW*0.5)**2 + (y - IH*0.5)**2)
    v = np.clip(1 - d/(IW*0.65), 0, 1) * 8
    for c in range(3):
        arr[:,:,c] = np.clip(arr[:,:,c].astype(int) + v.astype(int), 0, 255).astype(np.uint8)
    return Image.fromarray(arr)

def bg_gold_glow():
    y, x = _g()
    d = np.sqrt((x - IW*0.5)**2 + (y - IH*0.15)**2)
    g = np.clip(1 - d/(IW*0.42), 0, 1)**2 * 0.28
    r = (g*212).astype(np.uint8)
    gr= (g*175).astype(np.uint8)
    b = (g*55 ).astype(np.uint8)
    return Image.fromarray(np.stack([r, gr, b], 2))

def to_stream(img):
    s = io.BytesIO(); img.save(s, 'PNG'); s.seek(0); return s

print("Generating backgrounds…")
BG = {
    'red':   to_stream(bg_red_glow()),
    'dark':  to_stream(bg_dark_vignette()),
    'num':   to_stream(bg_number_glow()),
    'tech':  to_stream(bg_tech_grid()),
    'gold':  to_stream(bg_gold_glow()),
    'black': to_stream(Image.new('RGB', (IW, IH), (0, 0, 0))),
}
def gb(k): BG[k].seek(0); return BG[k]

# ── Transition helpers ────────────────────────────────────────────────────────
P_NS  = "http://schemas.openxmlformats.org/presentationml/2006/main"
P14   = "http://schemas.microsoft.com/office/powerpoint/2010/main"

def _rm_tr(slide):
    sld = slide._element
    ex  = sld.find(f'{{{P_NS}}}transition')
    if ex is not None: sld.remove(ex)

def morph(slide):
    _rm_tr(slide)
    slide._element.append(etree.fromstring(
        f'<p:transition xmlns:p="{P_NS}" xmlns:p14="{P14}" spd="med">'
        f'<p14:morph style="byObject"/></p:transition>'))

def fade_tr(slide):
    _rm_tr(slide)
    slide._element.append(etree.fromstring(
        f'<p:transition xmlns:p="{P_NS}" spd="med"><p:fade/></p:transition>'))

# ── Animation helpers ─────────────────────────────────────────────────────────
_gid = [0]
_aid = [500]
def _gid_next(): _gid[0]+=1; return _gid[0]
def _aid_next(): _aid[0]+=1; return _aid[0]

def animate_fade(slide, shape_ids):
    if not shape_ids: return
    sld = slide._element
    ex  = sld.find(f'{{{P_NS}}}timing')
    if ex is not None: sld.remove(ex)
    pars, blds = [], []
    for spid in shape_ids:
        gid = _gid_next()
        a,b,c,d,e,f,h = [_aid_next() for _ in range(7)]
        pars.append(f"""<p:par xmlns:p="{P_NS}">
  <p:cTn id="{a}" fill="hold">
    <p:stCondLst><p:cond delay="indefin"/></p:stCondLst>
    <p:childTnLst><p:par><p:cTn id="{b}" fill="hold">
      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
      <p:childTnLst><p:par><p:cTn id="{c}" presetID="10" presetClass="entr"
        presetSubtype="0" fill="hold" grpId="{gid}" nodeType="clickEffect">
        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
        <p:childTnLst>
          <p:set><p:cBhvr><p:cTn id="{d}" dur="1" fill="hold"/>
            <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
          </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>
          <p:animEffect transition="in" filter="fade">
            <p:cBhvr><p:cTn id="{e}" dur="600"/>
              <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            </p:cBhvr>
          </p:animEffect>
        </p:childTnLst></p:cTn></p:par>
      </p:childTnLst></p:cTn></p:par></p:childTnLst>
  </p:cTn></p:par>""")
        blds.append(f'<p:bldP xmlns:p="{P_NS}" spid="{spid}" grpId="{gid}"/>')
    ri, si = _aid_next(), _aid_next()
    sld.append(etree.fromstring(
        f'<p:timing xmlns:p="{P_NS}"><p:tnLst><p:par>'
        f'<p:cTn id="{ri}" dur="indefin" restart="whenNotActive" nodeType="tmRoot">'
        f'<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="{si}" dur="indefin" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(pars)}</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrevClick" delay="0">'
        f'<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
        f'<p:bldLst>{"".join(blds)}</p:bldLst></p:timing>'))

# ── Shape helpers ─────────────────────────────────────────────────────────────

def ns(bg_key='dark'):
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BLACK
    pic = s.shapes.add_picture(gb(bg_key), 0, 0, SLD_W, SLD_H)
    pic.name = f'_bg_{bg_key}'
    return s

def t(slide, text, x, y, w, h, sz=44, bold=True, col=WHITE,
      align=PP_ALIGN.CENTER, name=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    if name: tb.name = name
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = col; r.font.name = 'Arial'; r.font.italic = italic
    return tb

def lbl(s, text, col=RED):
    return t(s, text.upper(), 0, Inches(0.45), SLD_W, Inches(0.4),
             sz=13, bold=True, col=col)

def hl(s, y=Inches(3.8), col=RED):
    ln = s.shapes.add_shape(1, Inches(6.17), y, Inches(1), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = col; ln.line.fill.background()

def feat(s, label_, val, y, aids=None):
    cx = Inches(2.5)
    t1 = t(s, label_, cx,            y, Inches(4),    Inches(0.42),
           sz=20, bold=False, col=LGRAY, align=PP_ALIGN.LEFT)
    t2 = t(s, val,    cx+Inches(4),  y, Inches(4.33), Inches(0.42),
           sz=20, bold=True,  col=WHITE, align=PP_ALIGN.RIGHT)
    sep = s.shapes.add_shape(1, cx, y+Inches(0.4), Inches(8.33), Inches(0.01))
    sep.fill.solid(); sep.fill.fore_color.rgb = DGRAY; sep.line.fill.background()
    if aids is not None: aids += [t1.shape_id, t2.shape_id]

def num_slide(bg_key, num_str, nm, sz=155):
    """Big number only — first of a morph pair."""
    s = ns(bg_key)
    n = t(s, num_str, 0, Inches(1.4), SLD_W, Inches(4.8),
          sz=sz, bold=True, col=WHITE, name=nm)
    fade_tr(s)
    return s

def reveal_slide(bg_key, num_str, nm, spec_str, model_str, sz=120):
    """Number + spec + model — second of a morph pair (gets Morph transition)."""
    s = ns(bg_key)
    t(s, num_str,   0, Inches(0.9), SLD_W, Inches(4.0), sz=sz,   bold=True,  col=WHITE, name=nm)
    sp = t(s, spec_str,  0, Inches(4.9), SLD_W, Inches(0.9), sz=28,  bold=False, col=WHITE)
    mo = t(s, model_str, 0, Inches(5.7), SLD_W, Inches(0.7), sz=18,  bold=True,  col=RED)
    morph(s)
    return s


# ══════════════════════════════════════════════════════════════════════════════
print("Building slides…")

# ─── PART 1: OPENING ────────────────────────────────────────────────────────

# 1 — TCL
s = ns('red')
t(s, "TCL", 0, Inches(2.5), SLD_W, Inches(2.2), sz=120, bold=True, col=RED, name='tcl_logo')
fade_tr(s)

# 2 — TCL + 2026 (Morph: logo moves up, year appears)
s = ns('red')
t(s, "TCL",  0, Inches(0.9), SLD_W, Inches(1.2), sz=60, bold=True, col=RED, name='tcl_logo')
t(s, "2026", 0, Inches(1.8), SLD_W, Inches(4.8), sz=160, bold=True, col=WHITE, name='yr_2026')
morph(s)

# 3 — Opening statement
s = ns('red')
lbl(s, "Training Material")
t(s, "The Complete\nTCL TV Lineup.", 0, Inches(1.4), SLD_W, Inches(3.2), sz=72, bold=True, col=WHITE)
hl(s, Inches(4.55))
sub = t(s, "10 Models. One Vision.", 0, Inches(4.8), SLD_W, Inches(0.8), sz=26, bold=False, col=LGRAY)
animate_fade(s, [sub.shape_id])
fade_tr(s)

# 4 — Not a catalogue
s = ns('dark')
a1 = t(s, "This is not a product catalogue.", Inches(0.5), Inches(1.8), Inches(12.33), Inches(1.4), sz=52, bold=True,  col=WHITE)
a2 = t(s, "This is your competitive edge.",   Inches(0.5), Inches(3.4), Inches(12.33), Inches(1.2), sz=40, bold=False, col=LGRAY)
animate_fade(s, [a1.shape_id, a2.shape_id])
fade_tr(s)

# 5 — Let's begin
s = ns('dark')
t(s, "Let's begin.", 0, Inches(2.4), SLD_W, Inches(2.2), sz=90, bold=True, col=WHITE, name='lets_begin')
hl(s, Inches(4.55))
fade_tr(s)

# ─── PART 2: OVERVIEW ───────────────────────────────────────────────────────

# 6 — What changed
s = ns('dark')
lbl(s, "2026")
t(s, "What changed\nthis year?", 0, Inches(1.6), SLD_W, Inches(3.6), sz=72, bold=True, col=WHITE)
fade_tr(s)

# 7 — Everything (large number-style)
s = ns('num')
t(s, "Everything.", 0, Inches(1.5), SLD_W, Inches(4.2), sz=110, bold=True, col=WHITE, name='ev_big')
morph(s)

# 8 — Everything expanded (Morph: text shifts up, bullets appear)
s = ns('num')
t(s, "Everything.", 0, Inches(0.6), SLD_W, Inches(1.6), sz=56, bold=True, col=WHITE, name='ev_big')
hl(s, Inches(2.2))
b1 = t(s, "SQD Technology — across the premium range.",    Inches(1), Inches(2.5), Inches(11.33), Inches(0.7), sz=24, bold=False, col=LGRAY)
b2 = t(s, "RGB MiniLED — for the ultra-large category.",   Inches(1), Inches(3.2), Inches(11.33), Inches(0.7), sz=24, bold=False, col=LGRAY)
b3 = t(s, "Dolby Vision 2 Max — exclusive to TCL.",        Inches(1), Inches(3.9), Inches(11.33), Inches(0.7), sz=24, bold=False, col=LGRAY)
animate_fade(s, [b1.shape_id, b2.shape_id, b3.shape_id])
morph(s)

# 9 — Full lineup table
s = ns('dark')
lbl(s, "2026 Full Lineup")
rows = [
    ("X11L",              "Flagship · SQD MiniLED · 10,000 nits"),
    ("RM9L",              "Ultra-Large · RGB MiniLED · 9,000 nits"),
    ("C8L",               "Premium · SQD MiniLED · 6,000 nits"),
    ("C7L",               "Best Seller · SQD MiniLED · 3,000 nits"),
    ("RM7L",              "Large · RGB MiniLED · 2,000 nits"),
    ("C6L",               "Mid Premium · QLED"),
    ("P8L / P7L / P6L",   "Value · 4K Smart TV"),
    ("V6D / T6D",         "Entry · Smart TV"),
]
aids = []
for i, (lb, vl) in enumerate(rows):
    feat(s, lb, vl, Inches(1.05) + Inches(0.72)*i, aids)
animate_fade(s, aids)
fade_tr(s)

# 10 — Know the lineup
s = ns('dark')
t(s, "Know the lineup.", 0, Inches(1.8), SLD_W, Inches(1.6), sz=72, bold=True, col=WHITE)
t(s, "Own the floor.",   0, Inches(3.4), SLD_W, Inches(1.6), sz=72, bold=True, col=RED)
fade_tr(s)

# ─── PART 3: X11L ───────────────────────────────────────────────────────────

# 11 — X11L hero
s = ns('dark')
lbl(s, "Flagship")
t(s, "X11L", 0, Inches(0.8), SLD_W, Inches(5.0), sz=190, bold=True, col=WHITE, name='x11l_hero')
t(s, "The pinnacle of display technology.", 0, Inches(5.9), SLD_W, Inches(0.8), sz=22, bold=False, col=GRAY)
fade_tr(s)

# 12 — Statement
s = ns('dark')
t(s, "There is no TV\nabove the X11L.", 0, Inches(1.6), SLD_W, Inches(4.0), sz=80, bold=True, col=WHITE)
fade_tr(s)

# 13–14: 10,000 nits
num_slide(   'num', "10,000",  "n_x11l_nits", sz=155)
reveal_slide('num', "10,000",  "n_x11l_nits", "Peak Nits of Brightness", "TCL X11L")

# 15–16: 20,000+ zones
num_slide(   'num', "20,000+", "n_x11l_zones", sz=140)
reveal_slide('num', "20,000+", "n_x11l_zones", "Local Dimming Zones", "The most precise backlight ever built.")

# 17 — WHVA
s = ns('dark')
lbl(s, "Panel Technology")
t(s, "WHVA 2.0 Ultra", 0, Inches(1.4), SLD_W, Inches(2.0), sz=80, bold=True, col=WHITE)
hl(s, Inches(3.4))
w1 = t(s, "Wide High Viewing Angle — 2nd Generation.",      0, Inches(3.6), SLD_W, Inches(0.8), sz=26, bold=False, col=LGRAY)
w2 = t(s, "Perfect color from every seat in the room.",     0, Inches(4.4), SLD_W, Inches(0.8), sz=26, bold=False, col=LGRAY)
animate_fade(s, [w1.shape_id, w2.shape_id])
fade_tr(s)

# 18 — X11L specs table
s = ns('dark')
lbl(s, "X11L Full Specs")
aids = []
for lb, vl in [("Panel","WHVA 2.0 Ultra"),("Backlight","SQD MiniLED"),
                ("Peak Brightness","10,000 nits"),("Dimming Zones","20,000+"),
                ("HDR Support","Dolby Vision · HDR10+"),("Screen Sizes",'75"  ·  85"  ·  98"')]:
    feat(s, lb, vl, Inches(1.1)+Inches(0.82)*[0,1,2,3,4,5][["Panel","Backlight","Peak Brightness","Dimming Zones","HDR Support","Screen Sizes"].index(lb)], aids)
animate_fade(s, aids)
fade_tr(s)

# 19 — Sizes
s = ns('dark')
lbl(s, "Available In")
t(s, '75"     85"     98"', 0, Inches(2.0), SLD_W, Inches(2.5), sz=90, bold=True, col=WHITE)
hl(s, Inches(4.5))
t(s, "From premium to palace.", 0, Inches(4.7), SLD_W, Inches(0.8), sz=28, bold=False, col=LGRAY)
fade_tr(s)

# 20 — X11L close
s = ns('dark')
t(s, "There's nothing\nabove it.", 0, Inches(1.6), SLD_W, Inches(3.6), sz=88, bold=True, col=WHITE)
t(s, "TCL X11L", 0, Inches(5.3), SLD_W, Inches(0.7), sz=22, bold=True, col=RED)
fade_tr(s)

# ─── PART 4: RM9L ───────────────────────────────────────────────────────────

# 21 — RM9L hero
s = ns('dark')
lbl(s, "Ultra-Large Premium")
t(s, "RM9L", 0, Inches(0.8), SLD_W, Inches(5.0), sz=190, bold=True, col=WHITE)
t(s, "RGB MiniLED. Reimagined.", 0, Inches(5.9), SLD_W, Inches(0.8), sz=22, bold=False, col=GRAY)
fade_tr(s)

# 22 — Massive ×3
s = ns('dark')
m1 = t(s, "Massive screen.",    0, Inches(0.9), SLD_W, Inches(1.6), sz=72, bold=True, col=WHITE)
m2 = t(s, "Massive brightness.",0, Inches(2.5), SLD_W, Inches(1.6), sz=72, bold=True, col=WHITE)
m3 = t(s, "Massive impact.",    0, Inches(4.1), SLD_W, Inches(1.6), sz=72, bold=True, col=RED)
animate_fade(s, [m1.shape_id, m2.shape_id, m3.shape_id])
fade_tr(s)

# 23–24: 9,000 nits
num_slide(   'num', "9,000", "n_rm9l_nits", sz=155)
reveal_slide('num', "9,000", "n_rm9l_nits", "Peak Nits · RGB MiniLED", "TCL RM9L")

# 25 — Sizes
s = ns('dark')
lbl(s, "Screen Sizes")
t(s, '85"     98"     115"', 0, Inches(2.0), SLD_W, Inches(2.5), sz=80, bold=True, col=WHITE)
hl(s, Inches(4.5))
t(s, "When bigger is the point.", 0, Inches(4.7), SLD_W, Inches(0.8), sz=28, bold=False, col=LGRAY)
fade_tr(s)

# 26 — RM9L specs
s = ns('dark')
lbl(s, "RM9L Full Specs")
aids = []
for i, (lb, vl) in enumerate([("Backlight","RGB MiniLED"),("Peak Brightness","9,000 nits"),
                                ("Color Volume","Ultra-wide color gamut"),
                                ("HDR Support","Dolby Vision · HDR10+"),
                                ("Screen Sizes",'85"  ·  98"  ·  115"')]):
    feat(s, lb, vl, Inches(1.4)+Inches(0.88)*i, aids)
animate_fade(s, aids)
fade_tr(s)

# 27 — RM9L close
s = ns('dark')
t(s, "For the customer who wants it all\n— and larger.", Inches(0.5), Inches(1.6), Inches(12.33), Inches(4.2), sz=68, bold=True, col=WHITE)
fade_tr(s)

# ─── PART 5: C8L ────────────────────────────────────────────────────────────

# 28 — C8L hero
s = ns('dark')
lbl(s, "Premium")
t(s, "C8L", 0, Inches(0.7), SLD_W, Inches(5.2), sz=200, bold=True, col=WHITE)
t(s, "SQD MiniLED. Perfected.", 0, Inches(5.9), SLD_W, Inches(0.8), sz=22, bold=False, col=GRAY)
fade_tr(s)

# 29–30: 6,000 nits
num_slide(   'num', "6,000", "n_c8l_nits", sz=155)
reveal_slide('num', "6,000", "n_c8l_nits", "Peak Nits · SQD MiniLED", "TCL C8L")

# 31–32: 4,032 zones
num_slide(   'num', "4,032", "n_c8l_zones", sz=155)
reveal_slide('num', "4,032", "n_c8l_zones", "Local Dimming Zones", "Precision you can see.")

# 33 — Dolby Vision 2 Max EXCLUSIVE
s = ns('dark')
lbl(s, "World Exclusive")
t(s, "Dolby Vision 2 Max", 0, Inches(1.3), SLD_W, Inches(2.0), sz=76, bold=True, col=WHITE)
badge = s.shapes.add_shape(1, Inches(9.6), Inches(1.5), Inches(2.7), Inches(0.55))
badge.fill.solid(); badge.fill.fore_color.rgb = RED; badge.line.fill.background()
t(s, "EXCLUSIVE", Inches(9.6), Inches(1.5), Inches(2.7), Inches(0.55), sz=18, bold=True, col=WHITE)
hl(s, Inches(3.3))
e1 = t(s, "Only available on the TCL C8L.",         0, Inches(3.6), SLD_W, Inches(0.8), sz=28, bold=False, col=LGRAY)
e2 = t(s, "No other brand has this. Use it.",        0, Inches(4.4), SLD_W, Inches(0.8), sz=28, bold=False, col=RED)
animate_fade(s, [e1.shape_id, e2.shape_id])
fade_tr(s)

# 34 — C8L specs
s = ns('dark')
lbl(s, "C8L Full Specs")
aids = []
for i, (lb, vl) in enumerate([("Technology","SQD MiniLED"),("Peak Brightness","6,000 nits"),
                                ("Dimming Zones","4,032"),
                                ("Exclusive HDR","★  Dolby Vision 2 Max"),
                                ("HDR also","HDR10+"),("Screen Sizes",'55"  to  98"')]):
    feat(s, lb, vl, Inches(1.1)+Inches(0.82)*i, aids)
animate_fade(s, aids)
fade_tr(s)

# 35 — C8L close
s = ns('dark')
t(s, "Premium picture.",      0, Inches(1.6), SLD_W, Inches(1.6), sz=80, bold=True, col=WHITE)
t(s, "Exclusive technology.", 0, Inches(3.3), SLD_W, Inches(1.6), sz=80, bold=True, col=RED)
t(s, "TCL C8L", 0, Inches(5.3), SLD_W, Inches(0.7), sz=22, bold=True, col=RED)
fade_tr(s)

# ─── PART 6: C7L ────────────────────────────────────────────────────────────

# 36 — C7L hero
s = ns('dark')
lbl(s, "Best Seller")
t(s, "C7L", 0, Inches(0.7), SLD_W, Inches(5.2), sz=200, bold=True, col=WHITE)
t(s, "SQD MiniLED. For everyone.", 0, Inches(5.9), SLD_W, Inches(0.8), sz=22, bold=False, col=GRAY)
fade_tr(s)

# 37 — Statement
s = ns('dark')
t(s, "The one they'll\nask for by name.", 0, Inches(1.6), SLD_W, Inches(4.0), sz=80, bold=True, col=WHITE)
fade_tr(s)

# 38–39: 3,000 nits
num_slide(   'num', "3,000", "n_c7l_nits", sz=155)
reveal_slide('num', "3,000", "n_c7l_nits", "Peak Nits · SQD MiniLED", "TCL C7L")

# 40–41: 2,176 zones
num_slide(   'num', "2,176", "n_c7l_zones", sz=155)
reveal_slide('num', "2,176", "n_c7l_zones", "Local Dimming Zones", "Blacks this deep have never been this affordable.")

# 42–43: 144Hz
num_slide(   'num', "144Hz", "n_c7l_hz", sz=140)
reveal_slide('num', "144Hz", "n_c7l_hz", "Refresh Rate", "Zero blur. Zero ghosting.")

# 44 — Use cases
s = ns('dark')
u1 = t(s, "Gaming.",      0, Inches(0.6), SLD_W, Inches(1.6), sz=80, bold=True, col=WHITE)
u2 = t(s, "Movies.",      0, Inches(2.2), SLD_W, Inches(1.6), sz=80, bold=True, col=WHITE)
u3 = t(s, "Sports.",      0, Inches(3.8), SLD_W, Inches(1.6), sz=80, bold=True, col=WHITE)
u4 = t(s, "Everything.",  0, Inches(5.4), SLD_W, Inches(1.6), sz=80, bold=True, col=RED)
animate_fade(s, [u1.shape_id, u2.shape_id, u3.shape_id, u4.shape_id])
fade_tr(s)

# 45 — C7L specs
s = ns('dark')
lbl(s, "C7L Full Specs")
aids = []
for i, (lb, vl) in enumerate([("Technology","SQD MiniLED"),("Peak Brightness","3,000 nits"),
                                ("Dimming Zones","Up to 2,176"),("Refresh Rate","144Hz"),
                                ("HDR Support","Dolby Vision · HDR10+"),
                                ("Screen Sizes",'55"  to  98"'),("Color Combinations","1.07 Billion")]):
    feat(s, lb, vl, Inches(0.9)+Inches(0.72)*i, aids)
animate_fade(s, aids)
fade_tr(s)

# 46 — C7L close
s = ns('dark')
t(s, "The best seller\nfor a reason.", 0, Inches(1.6), SLD_W, Inches(3.6), sz=88, bold=True, col=WHITE)
t(s, "TCL C7L", 0, Inches(5.3), SLD_W, Inches(0.7), sz=22, bold=True, col=RED)
fade_tr(s)

# ─── PART 7: RM7L ───────────────────────────────────────────────────────────

# 47 — RM7L hero
s = ns('dark')
lbl(s, "Large Screen Value")
t(s, "RM7L", 0, Inches(0.8), SLD_W, Inches(5.0), sz=190, bold=True, col=WHITE)
t(s, "RGB MiniLED. Bold colors at scale.", 0, Inches(5.9), SLD_W, Inches(0.8), sz=22, bold=False, col=GRAY)
fade_tr(s)

# 48–49: 2,000 nits
num_slide(   'num', "2,000", "n_rm7l_nits", sz=155)
reveal_slide('num', "2,000", "n_rm7l_nits", "Peak Nits · RGB MiniLED", "TCL RM7L")

# 50 — RM7L close
s = ns('dark')
t(s, "Big screen.",   0, Inches(1.0), SLD_W, Inches(1.6), sz=80, bold=True, col=WHITE)
t(s, "Bold colors.",  0, Inches(2.6), SLD_W, Inches(1.6), sz=80, bold=True, col=WHITE)
t(s, "Better value.", 0, Inches(4.2), SLD_W, Inches(1.6), sz=80, bold=True, col=RED)
fade_tr(s)

# ─── PART 8: MID RANGE ──────────────────────────────────────────────────────

# 51 — Mid range section
s = ns('dark')
lbl(s, "Mid Range")
t(s, "C6L · P8L · P7L · P6L", 0, Inches(2.0), SLD_W, Inches(2.0), sz=60, bold=True, col=WHITE)
hl(s, Inches(4.0))
t(s, "Smart. Sharp. Affordable.", 0, Inches(4.3), SLD_W, Inches(0.8), sz=28, bold=False, col=LGRAY)
fade_tr(s)

# 52 — C6L
s = ns('dark')
lbl(s, "QLED · Mid Premium")
t(s, "C6L", 0, Inches(0.7), SLD_W, Inches(3.8), sz=160, bold=True, col=WHITE)
hl(s, Inches(4.5))
aids = []
for i, (lb, vl) in enumerate([("Technology","QLED"),("Resolution","4K Ultra HD"),("HDR","Dolby Vision · HDR10+")]):
    feat(s, lb, vl, Inches(4.8)+Inches(0.72)*i, aids)
animate_fade(s, aids)
fade_tr(s)

# 53 — P8L
s = ns('dark')
lbl(s, "MiniLED Entry")
t(s, "P8L", 0, Inches(0.7), SLD_W, Inches(3.8), sz=160, bold=True, col=WHITE)
hl(s, Inches(4.5))
aids = []
for i, (lb, vl) in enumerate([("Technology","MiniLED"),("Resolution","4K Ultra HD"),("HDR","HDR10+")]):
    feat(s, lb, vl, Inches(4.8)+Inches(0.72)*i, aids)
animate_fade(s, aids)
fade_tr(s)

# 54 — P7L + P6L
s = ns('dark')
lbl(s, "4K Smart TV · Value")
t(s, "P7L",  Inches(1.0), Inches(1.3), Inches(5.17), Inches(2.2), sz=110, bold=True, col=WHITE)
p1 = t(s, "4K · HDR10 · Smart TV\nGreat everyday viewing\nfor any family.",
       Inches(1.0), Inches(3.6), Inches(5.17), Inches(2.2), sz=22, bold=False, col=LGRAY)
vl_ = s.shapes.add_shape(1, Inches(6.62), Inches(1.3), Inches(0.02), Inches(5))
vl_.fill.solid(); vl_.fill.fore_color.rgb = DGRAY; vl_.line.fill.background()
t(s, "P6L",  Inches(7.2), Inches(1.3), Inches(5.17), Inches(2.2), sz=110, bold=True, col=WHITE)
p2 = t(s, "4K · HDR · Smart TV\nThe entry point\nto the TCL family.",
       Inches(7.2), Inches(3.6), Inches(5.17), Inches(2.2), sz=22, bold=False, col=LGRAY)
animate_fade(s, [p1.shape_id, p2.shape_id])
fade_tr(s)

# 55 — Mid-range close
s = ns('dark')
t(s, "A TCL for\nevery budget.", 0, Inches(1.4), SLD_W, Inches(3.6), sz=88, bold=True, col=WHITE)
hl(s, Inches(5.0))
c1 = t(s, "No customer walks away empty-handed.", 0, Inches(5.2), SLD_W, Inches(0.8), sz=26, bold=False, col=LGRAY)
animate_fade(s, [c1.shape_id])
fade_tr(s)

# ─── PART 9: ENTRY ──────────────────────────────────────────────────────────

# 56 — Entry section
s = ns('dark')
lbl(s, "Entry Level")
t(s, "V6D & T6D", 0, Inches(1.8), SLD_W, Inches(2.2), sz=88, bold=True, col=WHITE)
hl(s, Inches(4.0))
t(s, "Smart TV. Real quality. TCL price.", 0, Inches(4.3), SLD_W, Inches(0.8), sz=28, bold=False, col=LGRAY)
fade_tr(s)

# 57 — Entry specs
s = ns('dark')
lbl(s, "V6D & T6D Specs")
aids = []
for i, (lb, vl) in enumerate([("Resolution","Full HD / 4K"),("Smart OS","Google TV"),
                                ("HDR","HDR10"),("Connectivity","WiFi · Bluetooth · HDMI"),
                                ("Positioning","Best first TV")]):
    feat(s, lb, vl, Inches(1.6)+Inches(0.88)*i, aids)
animate_fade(s, aids)
fade_tr(s)

# 58 — Entry close
s = ns('dark')
t(s, "Everyone deserves", 0, Inches(1.8), SLD_W, Inches(1.6), sz=72, bold=True, col=WHITE)
t(s, "a great TV.",       0, Inches(3.5), SLD_W, Inches(1.6), sz=72, bold=True, col=RED)
fade_tr(s)

# ─── PART 10: SQD DEEP DIVE ─────────────────────────────────────────────────

# 59 — SQD hero
s = ns('tech')
lbl(s, "Technology Deep Dive")
t(s, "SQD", 0, Inches(0.6), SLD_W, Inches(5.6), sz=200, bold=True, col=RED, name='sqd_big')
fade_tr(s)

# 60 — SQD definition (Morph: SQD shrinks, definition appears)
s = ns('tech')
t(s, "SQD", 0, Inches(0.3), SLD_W, Inches(2.0), sz=80, bold=True, col=RED, name='sqd_big')
t(s, "Sub-pixel Quantum Dot.", 0, Inches(2.2), SLD_W, Inches(1.8), sz=72, bold=True, col=WHITE)
hl(s, Inches(4.0))
d1 = t(s, "The next evolution beyond QLED.  TCL-exclusive. Patent-protected.",
       0, Inches(4.3), SLD_W, Inches(0.8), sz=22, bold=False, col=LGRAY)
animate_fade(s, [d1.shape_id])
morph(s)

# 61 — Regular QLED
s = ns('tech')
t(s, "Regular QLED:",        0, Inches(1.6), SLD_W, Inches(1.4), sz=60, bold=True,  col=WHITE)
t(s, "One filter. Whole screen.", 0, Inches(3.1), SLD_W, Inches(1.4), sz=52, bold=False, col=GRAY)
fade_tr(s)

# 62 — SQD vs QLED
s = ns('tech')
t(s, "SQD:",          0, Inches(1.0), SLD_W, Inches(1.4), sz=72, bold=True,  col=RED)
t(s, "Every pixel.",  0, Inches(2.4), SLD_W, Inches(1.4), sz=72, bold=True,  col=WHITE)
t(s, "Its own filter.", 0, Inches(3.9), SLD_W, Inches(1.4), sz=72, bold=False, col=LGRAY)
fade_tr(s)

# 63–64: 1.07 Billion
num_slide(   'tech', "1.07B", "n_1b", sz=140)
reveal_slide('tech', "1.07B", "n_1b", "Color Combinations", "The human eye can see about 10 million.")

# 65 — SQD vs competition (3 columns)
s = ns('tech')
lbl(s, "SQD vs The Competition")
cols_ = [
    ("QLED",  GRAY, "Panel-level quantum dot filter.\nGood color. Limited precision.\nOne filter for millions of pixels."),
    ("OLED",  GRAY, "Self-emissive pixels.\nGreat black levels.\nLower brightness. Burn-in risk."),
    ("SQD ★", RED,  "Sub-pixel level filter.\nPeak brightness of MiniLED.\nColor precision per pixel."),
]
for i, (ttl, col_, desc_) in enumerate(cols_):
    x_ = Inches(0.5)+Inches(4.28)*i
    t(s, ttl,  x_, Inches(1.3), Inches(3.8), Inches(1.1), sz=44, bold=True, col=col_)
    t(s, desc_,x_, Inches(2.6), Inches(3.8), Inches(2.5), sz=20, bold=False, col=LGRAY)
    if i < 2:
        vl_ = s.shapes.add_shape(1, x_+Inches(3.9), Inches(1.3), Inches(0.02), Inches(4.5))
        vl_.fill.solid(); vl_.fill.fore_color.rgb = DGRAY; vl_.line.fill.background()
fade_tr(s)

# ─── PART 11: SALES TRAINING ────────────────────────────────────────────────

# 66 — Sales section
s = ns('dark')
lbl(s, "Sales Training")
t(s, "How to sell TCL\nin 2026.", 0, Inches(1.7), SLD_W, Inches(3.6), sz=80, bold=True, col=WHITE)
fade_tr(s)

# 67 — 3-step close
s = ns('dark')
lbl(s, "The 3-Step Close")
st = [
    '01   Ask: "What do you watch most?"',
    "02   Match: Show the right model.",
    "03   Demo: Let the screen close the sale.",
]
aids = []
for i, line in enumerate(st):
    tb = t(s, line, Inches(1.5), Inches(1.8)+Inches(1.6)*i, Inches(10.33), Inches(1.4),
           sz=32, bold=True, col=WHITE, align=PP_ALIGN.LEFT)
    aids.append(tb.shape_id)
animate_fade(s, aids)
fade_tr(s)

# 68 — Show don't tell
s = ns('dark')
t(s, "Show.",       0, Inches(1.5), SLD_W, Inches(1.8), sz=100, bold=True, col=WHITE)
t(s, "Don't tell.", 0, Inches(3.3), SLD_W, Inches(1.8), sz=100, bold=True, col=RED)
c2 = t(s, "The SQD difference is visible in 5 seconds.", 0, Inches(5.4), SLD_W, Inches(0.8), sz=24, bold=False, col=LGRAY)
animate_fade(s, [c2.shape_id])
fade_tr(s)

# 69 — Demo script
s = ns('dark')
lbl(s, "The SQD Demo — In-Store Script")
script = [
    "1.  Point to any QLED TV next to C7L / C8L.",
    '2.  Say: "See the green on that screen?"',
    '3.  Point to TCL: "Now look at ours."',
    '4.  Say: "That\'s SQD — sub-pixel filtering."',
    '5.  Ask: "Can you see the difference?"',
]
aids = []
for i, ln in enumerate(script):
    tb = t(s, ln, Inches(1.5), Inches(1.4)+Inches(1.1)*i, Inches(10.33), Inches(0.9),
           sz=26, bold=False, col=WHITE, align=PP_ALIGN.LEFT)
    aids.append(tb.shape_id)
animate_fade(s, aids)
fade_tr(s)

# 70 — Upsell path
s = ns('dark')
lbl(s, "Upsell Path")
aids = []
for i, (lb, vl) in enumerate([
    ("Budget-conscious",         "C7L — Best value SQD"),
    ("Wants the best picture",   "C8L — Dolby Vision 2 Max"),
    ("Wants extreme size",       'RM9L — 85" to 115"'),
    ("Money is no object",       "X11L — The pinnacle"),
    ("Gaming focus",             "C7L — 144Hz SQD"),
    ("First-time buyer",         "P6L or P7L — Smart 4K"),
]):
    feat(s, lb, vl, Inches(1.1)+Inches(0.82)*i, aids)
animate_fade(s, aids)
fade_tr(s)

# ─── PART 12: CLOSING ───────────────────────────────────────────────────────

# 71 — You are the expert
s = ns('dark')
t(s, "You are now\nthe expert.", 0, Inches(1.6), SLD_W, Inches(3.6), sz=88, bold=True, col=WHITE)
fade_tr(s)

# 72 — TCL CERTIFIED (gold glow background)
s = ns('gold')
t(s, "ACHIEVEMENT", 0, Inches(0.7), SLD_W, Inches(0.5), sz=13, bold=True, col=GOLD)
t(s, "TCL\nCERTIFIED.", 0, Inches(1.0), SLD_W, Inches(4.0), sz=130, bold=True, col=GOLD)
hl(s, Inches(5.1), col=GOLD)
c3 = t(s, "You now sell the best TV technology in the world.",
       0, Inches(5.4), SLD_W, Inches(0.8), sz=24, bold=False, col=LGRAY)
t(s, "TCL  ·  2026", 0, Inches(6.5), SLD_W, Inches(0.7), sz=16, bold=True, col=RED)
animate_fade(s, [c3.shape_id])
fade_tr(s)

# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"\n✅  Saved: {OUTPUT}")
print(f"   Slides : {len(prs.slides)}")
print(f"   Size   : {os.path.getsize(OUTPUT):,} bytes")
